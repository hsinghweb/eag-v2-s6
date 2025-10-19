import os
from mcp.server.fastmcp import FastMCP, Image
from mcp.server.fastmcp.prompts import base
from mcp.types import TextContent
from PIL import Image as PILImage
import math
import sys
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.util import Pt
import asyncio
from tools import (
    number_list_to_sum,
    calculate_difference,
    number_list_to_product,
    calculate_division,
    strings_to_chars_to_int as local_strings_to_chars_to_int,
    int_list_to_exponential_values,
    fibonacci_numbers as local_fibonacci_numbers,
    calculate_factorial,
    calculate_permutation,
    calculate_combination,
    calculate_salary_for_id,
    calculate_salary_for_name,
    calculate_percentage,
)
import logging
from datetime import datetime
import traceback
import smtplib
from email.mime.text import MIMEText
from dotenv import load_dotenv

# Load environment variables
load_dotenv()

# Configure logging
log_dir = "logs"
os.makedirs(log_dir, exist_ok=True)
log_file = os.path.join(log_dir, f"mcp_server_{datetime.now().strftime('%Y%m%d_%H%M%S')}.log")

# Define constant for PowerPoint filename
PPTX_FILENAME = 'presentation.pptx'

logging.basicConfig(
    level=logging.DEBUG,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler(log_file),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

# Instantiate an MCP server client
mcp = FastMCP("Calculator")

# DEFINE TOOLS

# Addition tool
@mcp.tool()
def add(a: int, b: int) -> int:
    """Add two numbers"""
    logger.info(f"Calling add(a: {a}, b: {b}) -> int")
    return int(a + b)

@mcp.tool()
def add_list(lst: list) -> int:
    """Add all numbers in a list"""
    logger.info(f"Calling add_list(lst: {lst}) -> int")
    return sum(lst)

# Subtraction tool
@mcp.tool()
def subtract(a: int, b: int) -> int:
    """Subtract two numbers"""
    logger.info(f"Calling subtract(a: {a}, b: {b}) -> int")
    return int(a - b)

# Multiplication tool
@mcp.tool()
def multiply(a: int, b: int) -> int:
    """Multiply two numbers"""
    logger.info(f"Calling multiply(a: {a}, b: {b}) -> int")
    return int(a * b)

# Division tool
@mcp.tool() 
def divide(a: int, b: int) -> float:
    """Divide two numbers"""
    logger.info(f"Calling divide(a: {a}, b: {b}) -> float")
    return float(a / b)

# Power tool
@mcp.tool()
def power(a: int, b: int) -> int:
    """Power of two numbers"""
    logger.info(f"Calling power(a: {a}, b: {b}) -> int")
    return int(a ** b)

# Square root tool
@mcp.tool()
def sqrt(a: int) -> float:
    """Square root of a number"""
    logger.info(f"Calling sqrt(a: {a}) -> float")
    return float(a ** 0.5)

# Cube root tool
@mcp.tool()
def cbrt(a: int) -> float:
    """Cube root of a number"""
    logger.info(f"Calling cbrt(a: {a}) -> float")
    return float(a ** (1/3))

# Factorial tool
@mcp.tool()
def factorial(a: int) -> int:
    """Factorial of a number"""
    logger.info(f"Calling factorial(a: {a}) -> int")
    return int(math.factorial(a))

@mcp.tool()
def create_thumbnail(image_path: str) -> Image:
    """Create a thumbnail from an image"""
    logger.info(f"Calling create_thumbnail(image_path: {image_path}) -> Image")
    img = PILImage.open(image_path)
    img.thumbnail((100, 100))
    return Image(data=img.tobytes(), format="png")

@mcp.tool()
def strings_to_chars_to_int(string: str) -> list[int]:
    """Return the ASCII values of the characters in a word"""
    logger.info(f"Calling strings_to_chars_to_int(string: {string}) -> list[int]")
    return [int(ord(char)) for char in string]

@mcp.tool()
def int_list_to_exponential_sum(int_list: list) -> float:
    """Return sum of exponentials of numbers in a list"""
    logger.info(f"Calling int_list_to_exponential_sum(int_list: {int_list}) -> float")
    return sum(math.exp(i) for i in int_list)

@mcp.tool()
def fibonacci_numbers(n: int) -> list:
    """Return the first n Fibonacci Numbers"""
    logger.info(f"Calling fibonacci_numbers(n: {n}) -> list")
    if n <= 0:
        return []
    fib_sequence = [0, 1]
    for _ in range(2, n):
        fib_sequence.append(fib_sequence[-1] + fib_sequence[-2])
    return fib_sequence[:n]

@mcp.tool()
async def close_powerpoint() -> dict:
    """Close PowerPoint"""
    try:
        logger.info("Calling close_powerpoint()")
        # Suppress all output from taskkill
        try:
            # Use subprocess with output suppressed
            proc = await asyncio.create_subprocess_shell(
                'taskkill /F /IM POWERPNT.EXE',
                stdout=asyncio.subprocess.DEVNULL,
                stderr=asyncio.subprocess.DEVNULL
            )
            await proc.communicate()
        except Exception as e:
            logger.warning(f"taskkill failed or PowerPoint not running: {e}")
        await asyncio.sleep(2)
        logger.info("PowerPoint closed successfully")
        return {
            "content": [
                TextContent(
                    type="text",
                    text="PowerPoint closed successfully"
                )
            ]
        }
    except Exception as e:
        logger.error(f"Error in close_powerpoint: {str(e)}")
        logger.error(traceback.format_exc())
        return {
            "content": [
                TextContent(
                    type="text",
                    text=f"Error closing PowerPoint: {str(e)}"
                )
            ]
        }

@mcp.tool()
async def open_powerpoint() -> dict:
    """Open a new PowerPoint presentation"""
    try:
        logger.info("Calling open_powerpoint()")
        await close_powerpoint()
        await asyncio.sleep(3)
        prs = Presentation()
        prs.slide_layouts[0]
        filename = PPTX_FILENAME
        prs.save(filename)
        await asyncio.sleep(5)
        # Open the presentation, suppressing output
        try:
            # On Windows, os.startfile does not print to stdout, but just in case:
            # Use subprocess with output suppressed for other OSes if needed
            os.startfile(filename)
        except Exception as e:
            logger.warning(f"os.startfile failed: {e}")
        await asyncio.sleep(10)
        logger.info("PowerPoint opened successfully with a new presentation")
        return {
            "content": [
                TextContent(
                    type="text",
                    text="PowerPoint opened successfully with a new presentation"
                )
            ]
        }
    except Exception as e:
        logger.error(f"Error in open_powerpoint: {str(e)}")
        logger.error(traceback.format_exc())
        return {
            "content": [
                TextContent(
                    type="text",
                    text=f"Error opening PowerPoint: {str(e)}"
                )
            ]
        }

@mcp.tool()
async def draw_rectangle(x1: int = 1, y1: int = 1, x2: int = 8, y2: int = 6) -> dict:
    """Draw a rectangle in the first slide of PowerPoint
    
    Args:
        x1: X-coordinate of top-left corner (1-8, default: 1)
        y1: Y-coordinate of top-left corner (1-8, default: 1)
        x2: X-coordinate of bottom-right corner (1-8, default: 8)
        y2: Y-coordinate of bottom-right corner (1-8, default: 6)
    """
    try:
        logger.info(f"Drawing rectangle with parameters: x1={x1} ({type(x1)}), y1={y1} ({type(y1)}), x2={x2} ({type(x2)}), y2={y2} ({type(y2)})")
        
        # Convert parameters to integers
        try:
            x1 = int(float(str(x1)))
            y1 = int(float(str(y1)))
            x2 = int(float(str(x2)))
            y2 = int(float(str(y2)))
        except (ValueError, TypeError) as e:
            error_msg = f"Failed to convert parameters to integers: {str(e)}"
            logger.error(error_msg)
            return {"content": [TextContent(type="text", text=error_msg)]}

        logger.debug(f"Converted coordinates: ({x1},{y1}) to ({x2},{y2})")
        
        # Validate coordinates
        if not (1 <= x1 <= 8 and 1 <= y1 <= 8 and 1 <= x2 <= 8 and 1 <= y2 <= 8):
            error_msg = f"Coordinates must be between 1 and 8, got: ({x1},{y1}) to ({x2},{y2})"
            logger.error(error_msg)
            return {"content": [TextContent(type="text", text=error_msg)]}
        
        if x2 <= x1 or y2 <= y1:
            error_msg = f"End coordinates must be greater than start coordinates: ({x1},{y1}) to ({x2},{y2})"
            logger.error(error_msg)
            return {"content": [TextContent(type="text", text=error_msg)]}
        
        # Wait before modifying the presentation
        await asyncio.sleep(2)
        
        # Ensure PowerPoint is closed before modifying the file
        await close_powerpoint()
        await asyncio.sleep(2)
        prs = Presentation(PPTX_FILENAME)
        slide = prs.slides[0]
        
        # Store existing text boxes
        slide = prs.slides[0]
        
        # Store existing text boxes
        text_boxes = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height
                text_boxes.append((text, left, top, width, height))
        
        # Clear existing shapes except text boxes
        for shape in slide.shapes:
            if not shape.has_text_frame:
                sp = shape._element
                sp.getparent().remove(sp)
        
        # Convert coordinates to inches
        left = Inches(x1)
        top = Inches(y1)
        width = Inches(x2 - x1)
        height = Inches(y2 - y1)
        
        logger.debug(f"Rectangle dimensions - left={left}, top={top}, width={width}, height={height}")
        
        # Add rectangle
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            left, top, width, height
        )
        
        # Make the rectangle more visible
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White fill
        shape.line.color.rgb = RGBColor(0, 0, 0)  # Black border
        prs.save(PPTX_FILENAME)
        await asyncio.sleep(2)
        
        # Reopen PowerPoint
        os.startfile(PPTX_FILENAME)
        await asyncio.sleep(5)
        # Reopen PowerPoint
        os.startfile(PPTX_FILENAME)
        await asyncio.sleep(5)
        
        logger.info(f"Rectangle drawn successfully from ({x1},{y1}) to ({x2},{y2})")
        return {
            "content": [
                TextContent(
                    type="text",
                    text=f"Rectangle drawn successfully from ({x1},{y1}) to ({x2},{y2})"
                )
            ]
        }
            
    except Exception as e:
        error_msg = f"Error in draw_rectangle: {str(e)}"
        logger.error(error_msg)
        logger.error(traceback.format_exc())
        return {"content": [TextContent(type="text", text=error_msg)]}

@mcp.tool()
async def add_text_in_powerpoint(text: str) -> dict:
    """Add text to the first slide of PowerPoint"""
    try:
        logger.info(f"Received text to add: {text}")
        logger.debug(f"Text type: {type(text)}")
        logger.debug(f"Text length: {len(text)}")
        logger.debug(f"Text contains newlines: {'\\n' in text}")
        
        # Wait before adding text
        await asyncio.sleep(5)
        
        # Ensure PowerPoint is closed before modifying the file
        await close_powerpoint()
        prs = Presentation(PPTX_FILENAME)
        slide = prs.slides[0]
        
        # Add a text box positioned inside the rectangle
        slide = prs.slides[0]
        
        # Add a text box positioned inside the rectangle
        # Match the rectangle position from draw_rectangle
        left = Inches(2.2)  # Slightly more than rectangle left for margin
        top = Inches(2.5)   # Centered vertically in rectangle
        width = Inches(4.6) # Slightly less than rectangle width for margin
        height = Inches(2)  # Enough height for text
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.clear()  # Clear existing text
        text_frame.word_wrap = True  # Enable word wrap
        text_frame.vertical_anchor = 1  # Middle vertical alignment
        
        # Split text into lines
        lines = text.split('\n')
        logger.debug(f"Number of lines: {len(lines)}")
        logger.debug(f"Lines to add: {lines}")
        
        # Add each line as a separate paragraph
        for i, line in enumerate(lines):
            if line.strip():  # Only add non-empty lines
                p = text_frame.add_paragraph()
                p.text = line.strip()
                p.alignment = 1  # Center align the text
                
                # Format the text
                run = p.runs[0]
                if "Final Result:" in line:
                    run.font.size = Pt(32)  # Header size
                    run.font.bold = True
                else:
                    run.font.size = Pt(28)  # Value size
                    run.font.bold = True
                
                run.font.color.rgb = RGBColor(0, 0, 0)  # Black text
        prs.save(PPTX_FILENAME)
        await asyncio.sleep(5)
        
        # Reopen PowerPoint
        os.startfile(PPTX_FILENAME)
        await asyncio.sleep(10)
        # Reopen PowerPoint
        os.startfile(PPTX_FILENAME)
        await asyncio.sleep(10)
        
        logger.info(f"Text added successfully: {text}")
        return {
            "content": [
                TextContent(
                    type="text",
                    text=f"Text added successfully: {text}"
                )
            ]
        }
    except Exception as e:
        logger.error(f"Error in add_text_in_powerpoint: {str(e)}")
        logger.error(traceback.format_exc())
        return {
            "content": [
                TextContent(
                    type="text",
                    text=f"Error adding text: {str(e)}"
                )
            ]
        }

@mcp.tool()
async def send_gmail(content: str) -> dict:
    """Send an email with the specified content via Gmail"""
    try:
        logger.info(f"Calling send_gmail(content: {content[:50]}...)")
        
        # Retrieve Gmail credentials and recipient from .env
        gmail_address = os.getenv("GMAIL_ADDRESS")
        gmail_app_password = os.getenv("GMAIL_APP_PASSWORD")
        recipient_email = os.getenv("RECIPIENT_EMAIL")
        
        if not all([gmail_address, gmail_app_password, recipient_email]):
            error_msg = "Missing GMAIL_ADDRESS, GMAIL_APP_PASSWORD, or RECIPIENT_EMAIL in .env file"
            logger.error(error_msg)
            return {
                "content": [
                    TextContent(
                        type="text",
                        text=error_msg
                    )
                ]
            }
        
        # Validate Gmail address and recipient email format
        if not (gmail_address.endswith('@gmail.com') and '@' in recipient_email):
            error_msg = f"Invalid email format: GMAIL_ADDRESS={gmail_address}, RECIPIENT_EMAIL={recipient_email}"
            logger.error(error_msg)
            return {
                "content": [
                    TextContent(
                        type="text",
                        text=error_msg
                    )
                ]
            }
        
        # Create the email message
        msg = MIMEText(content)
        msg['Subject'] = 'Math Agent Result'
        msg['From'] = gmail_address
        msg['To'] = recipient_email
        
        # Connect to Gmail's SMTP server
        try:
            with smtplib.SMTP_SSL('smtp.gmail.com', 465) as server:
                logger.debug("Connecting to Gmail SMTP server (smtp.gmail.com:465)")
                server.login(gmail_address, gmail_app_password)
                logger.debug(f"Logged in as {gmail_address}")
                server.sendmail(gmail_address, recipient_email, msg.as_string())
                logger.info(f"Email sent successfully to {recipient_email}")
        except smtplib.SMTPAuthenticationError as e:
            error_msg = f"SMTP Authentication failed: {str(e)}. Ensure GMAIL_APP_PASSWORD is correct and 2-Step Verification is enabled."
            logger.error(error_msg)
            return {
                "content": [
                    TextContent(
                        type="text",
                        text=error_msg
                    )
                ]
            }
        except Exception as e:
            error_msg = f"Failed to send email: {str(e)}"
            logger.error(error_msg)
            return {
                "content": [
                    TextContent(
                        type="text",
                        text=error_msg
                    )
                ]
            }
        
        return {
            "content": [
                TextContent(
                    type="text",
                    text=f"Email sent successfully to {recipient_email}"
                )
            ]
        }
    except Exception as e:
        error_msg = f"Error in send_gmail: {str(e)}"
        logger.error(error_msg)
        logger.error(traceback.format_exc())
        return {
            "content": [
                TextContent(
                    type="text",
                    text=error_msg
                )
            ]
        }

# DEFINE RESOURCES

# Add a dynamic greeting resource
@mcp.resource("greeting://{name}")
def get_greeting(name: str) -> str:
    """Get a personalized greeting"""
    logger.info(f"Calling get_greeting(name: {name}) -> str")
    return f"Hello, {name}!"

# DEFINE AVAILABLE PROMPTS
@mcp.prompt()
def review_code(code: str) -> str:
    logger.info(f"Calling review_code(code: {code[:50]}...) -> str")
    return f"Please review this code:\n\n{code}"

@mcp.prompt()
def debug_error(error: str) -> list[base.Message]:
    logger.info(f"Calling debug_error(error: {error[:50]}...) -> list[base.Message]")
    return [
        base.UserMessage("I'm seeing this error:"),
        base.UserMessage(error),
        base.AssistantMessage("I'll help debug that. What have you tried so far?"),
    ]

# TOOLS WRAPPING FUNCTIONS FROM tools.py
@mcp.tool()
def t_number_list_to_sum(lst: list) -> int:
    """Sum numbers in a list (wrapper around tools.number_list_to_sum)"""
    logger.info(f"Calling t_number_list_to_sum(lst: {lst}) -> int")
    return number_list_to_sum(lst)

@mcp.tool()
def t_calculate_difference(a: float, b: float) -> float:
    """Difference between two numbers (wrapper)"""
    logger.info(f"Calling t_calculate_difference(a: {a}, b: {b}) -> float")
    return calculate_difference(a, b)

@mcp.tool()
def t_number_list_to_product(lst: list) -> int:
    """Product of numbers in a list (wrapper)"""
    logger.info(f"Calling t_number_list_to_product(lst: {lst}) -> int")
    return number_list_to_product(lst)

@mcp.tool()
def t_calculate_division(a: float, b: float) -> float:
    """Division of two numbers (wrapper)"""
    logger.info(f"Calling t_calculate_division(a: {a}, b: {b}) -> float")
    return calculate_division(a, b)

@mcp.tool()
def t_strings_to_chars_to_int(s: str) -> list[int]:
    """ASCII values of characters (wrapper)"""
    logger.info(f"Calling t_strings_to_chars_to_int(s: {s}) -> list[int]")
    return local_strings_to_chars_to_int(s)

@mcp.tool()
def t_int_list_to_exponential_values(lst: list) -> list[float]:
    """Exponential of list elements (wrapper)"""
    logger.info(f"Calling t_int_list_to_exponential_values(lst: {lst}) -> list[float]")
    return int_list_to_exponential_values(lst)

@mcp.tool()
def t_fibonacci_numbers(n: int) -> list[int]:
    """First n Fibonacci numbers (wrapper)"""
    logger.info(f"Calling t_fibonacci_numbers(n: {n}) -> list[int]")
    return local_fibonacci_numbers(n)

@mcp.tool()
def t_calculate_factorial(n: int) -> list[int]:
    """List of factorials up to n-1 (wrapper)"""
    logger.info(f"Calling t_calculate_factorial(n: {n}) -> list[int]")
    return calculate_factorial(n)

@mcp.tool()
def t_calculate_permutation(n: int, r: int) -> int:
    """Permutation nPr (wrapper)"""
    logger.info(f"Calling t_calculate_permutation(n: {n}, r: {r}) -> int")
    return calculate_permutation(n, r)

@mcp.tool()
def t_calculate_combination(n: int, r: int) -> int:
    """Combination nCr (wrapper)"""
    logger.info(f"Calling t_calculate_combination(n: {n}, r: {r}) -> int")
    return calculate_combination(n, r)

@mcp.tool()
def t_calculate_salary_for_id(emp_id: int) -> float | int | None:
    """Salary by employee id (wrapper)"""
    logger.info(f"Calling t_calculate_salary_for_id(emp_id: {emp_id}) -> float | int | None")
    return calculate_salary_for_id(emp_id)

@mcp.tool()
def t_calculate_salary_for_name(emp_name: str) -> float | int | None:
    """Salary by employee name (wrapper)"""
    logger.info(f"Calling t_calculate_salary_for_name(emp_name: {emp_name}) -> float | int | None")
    return calculate_salary_for_name(emp_name)

@mcp.tool()
def t_calculate_percentage(percent: float, number: float) -> float:
    """Calculate percentage of a number (wrapper)"""
    logger.info(f"Calling t_calculate_percentage(percent: {percent}, number: {number}) -> float")
    return calculate_percentage(percent, number)

@mcp.tool()
def fallback_reasoning(description: str) -> str:
    """Fallback reasoning step when the agent is uncertain or a tool fails."""
    logger.info(f"Calling fallback_reasoning(description: {description})")
    return f"Fallback invoked: {description}"

if __name__ == "__main__":
    logger.info("Starting the MCP server")
    if len(sys.argv) > 1 and sys.argv[1] == "dev":
        logger.info("Running MCP server in dev mode without transport")
        mcp.run()  # Run without transport for dev server
    else:
        logger.info("Running MCP server with stdio transport")
        mcp.run(transport="stdio")  # Run with stdio for direct execution