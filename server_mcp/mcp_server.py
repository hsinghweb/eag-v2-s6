import os
import sys
from pathlib import Path

# Add parent directory to path to allow imports when run directly
sys.path.insert(0, str(Path(__file__).parent.parent))

from mcp.server.fastmcp import FastMCP, Image
from mcp.server.fastmcp.prompts import base
from mcp.types import TextContent
from PIL import Image as PILImage
import math
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.util import Pt
import asyncio
from server_mcp.tools import (
    number_list_to_sum,
    calculate_difference,
    number_list_to_product,
    calculate_division,
    strings_to_chars_to_int as local_strings_to_chars_to_int,
    int_list_to_exponential_values,
    fibonacci_numbers as local_fibonacci_numbers,
    calculate_factorial,
    calculate_permutation,
    calculate_combination,
    calculate_salary_for_id,
    calculate_salary_for_name,
    calculate_percentage,
)
from server_mcp.models import (
    DrawRectangleInput, DrawRectangleOutput,
    AddTextInput, AddTextOutput,
    PowerPointOperationOutput,
    SendGmailInput, SendGmailOutput,
    NumberListInput, NumberListOutput,
    TwoNumberInput, TwoNumberOutput,
    PercentageInput, PercentageOutput,
    StringToCharsInput, StringToCharsOutput,
    ExponentialInput, ExponentialOutput,
    FibonacciInput, FibonacciOutput,
    FactorialInput, FactorialOutput,
    PermutationInput, PermutationOutput,
    CombinationInput, CombinationOutput,
    EmployeeIdInput, EmployeeNameInput, SalaryOutput,
    FallbackInput, FallbackOutput,
)
import logging
from datetime import datetime
import traceback
import smtplib
from email.mime.text import MIMEText
from dotenv import load_dotenv

# Load environment variables
load_dotenv()

# Configure logging
log_dir = "logs"
os.makedirs(log_dir, exist_ok=True)
log_file = os.path.join(log_dir, f"mcp_server_{datetime.now().strftime('%Y%m%d_%H%M%S')}.log")

# Define constant for PowerPoint filename
PPTX_FILENAME = 'presentation.pptx'

logging.basicConfig(
    level=logging.DEBUG,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler(log_file),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

# Instantiate an MCP server client
mcp = FastMCP("Calculator")

# DEFINE TOOLS
@mcp.tool()
async def close_powerpoint() -> dict:
    """Close PowerPoint"""
    try:
        logger.info("Calling close_powerpoint()")
        # Suppress all output from taskkill
        try:
            # Use subprocess with output suppressed
            proc = await asyncio.create_subprocess_shell(
                'taskkill /F /IM POWERPNT.EXE',
                stdout=asyncio.subprocess.DEVNULL,
                stderr=asyncio.subprocess.DEVNULL
            )
            await proc.communicate()
        except Exception as e:
            logger.warning(f"taskkill failed or PowerPoint not running: {e}")
        await asyncio.sleep(2)
        logger.info("PowerPoint closed successfully")
        return {
            "content": [
                TextContent(
                    type="text",
                    text="PowerPoint closed successfully"
                )
            ]
        }
    except Exception as e:
        logger.error(f"Error in close_powerpoint: {str(e)}")
        logger.error(traceback.format_exc())
        return {
            "content": [
                TextContent(
                    type="text",
                    text=f"Error closing PowerPoint: {str(e)}"
                )
            ]
        }

@mcp.tool()
async def open_powerpoint() -> dict:
    """Open a new PowerPoint presentation"""
    try:
        logger.info("Calling open_powerpoint()")
        await close_powerpoint()
        await asyncio.sleep(3)
        prs = Presentation()
        prs.slide_layouts[0]
        filename = PPTX_FILENAME
        prs.save(filename)
        await asyncio.sleep(5)
        # Open the presentation, suppressing output
        try:
            # On Windows, os.startfile does not print to stdout, but just in case:
            # Use subprocess with output suppressed for other OSes if needed
            os.startfile(filename)
        except Exception as e:
            logger.warning(f"os.startfile failed: {e}")
        await asyncio.sleep(10)
        logger.info("PowerPoint opened successfully with a new presentation")
        return {
            "content": [
                TextContent(
                    type="text",
                    text="PowerPoint opened successfully with a new presentation"
                )
            ]
        }
    except Exception as e:
        logger.error(f"Error in open_powerpoint: {str(e)}")
        logger.error(traceback.format_exc())
        return {
            "content": [
                TextContent(
                    type="text",
                    text=f"Error opening PowerPoint: {str(e)}"
                )
            ]
        }

@mcp.tool()
async def draw_rectangle(input: DrawRectangleInput) -> dict:
    """Draw a rectangle in the first slide of PowerPoint"""
    try:
        logger.info(f"Drawing rectangle with validated parameters: ({input.x1},{input.y1}) to ({input.x2},{input.y2})")
        
        # Wait before modifying the presentation
        await asyncio.sleep(2)
        
        # Ensure PowerPoint is closed before modifying the file
        await close_powerpoint()
        await asyncio.sleep(2)
        prs = Presentation(PPTX_FILENAME)
        slide = prs.slides[0]
        
        # Store existing text boxes
        slide = prs.slides[0]
        
        # Store existing text boxes
        text_boxes = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height
                text_boxes.append((text, left, top, width, height))
        
        # Clear existing shapes except text boxes
        for shape in slide.shapes:
            if not shape.has_text_frame:
                sp = shape._element
                sp.getparent().remove(sp)
        
        # Convert coordinates to inches
        left = Inches(input.x1)
        top = Inches(input.y1)
        width = Inches(input.x2 - input.x1)
        height = Inches(input.y2 - input.y1)
        
        logger.debug(f"Rectangle dimensions - left={left}, top={top}, width={width}, height={height}")
        
        # Add rectangle
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            left, top, width, height
        )
        
        # Make the rectangle more visible
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White fill
        shape.line.color.rgb = RGBColor(0, 0, 0)  # Black border
        prs.save(PPTX_FILENAME)
        await asyncio.sleep(2)
        
        # Reopen PowerPoint
        os.startfile(PPTX_FILENAME)
        await asyncio.sleep(5)
        # Reopen PowerPoint
        os.startfile(PPTX_FILENAME)
        await asyncio.sleep(5)
        
        logger.info(f"Rectangle drawn successfully from ({input.x1},{input.y1}) to ({input.x2},{input.y2})")
        return {
            "content": [
                TextContent(
                    type="text",
                    text=f"Rectangle drawn successfully from ({input.x1},{input.y1}) to ({input.x2},{input.y2})"
                )
            ]
        }
            
    except Exception as e:
        error_msg = f"Error in draw_rectangle: {str(e)}"
        logger.error(error_msg)
        logger.error(traceback.format_exc())
        return {"content": [TextContent(type="text", text=error_msg)]}

@mcp.tool()
async def add_text_in_powerpoint(input: AddTextInput) -> dict:
    """Add text to the first slide of PowerPoint"""
    try:
        logger.info(f"Received text to add: {input.text}")
        logger.debug(f"Text length: {len(input.text)}")
        logger.debug(f"Text contains newlines: {'\\n' in input.text}")
        
        # Wait before adding text
        await asyncio.sleep(5)
        
        # Ensure PowerPoint is closed before modifying the file
        await close_powerpoint()
        prs = Presentation(PPTX_FILENAME)
        slide = prs.slides[0]
        
        # Add a text box positioned inside the rectangle
        slide = prs.slides[0]
        
        # Add a text box positioned inside the rectangle
        # Match the rectangle position from draw_rectangle
        left = Inches(2.2)  # Slightly more than rectangle left for margin
        top = Inches(2.5)   # Centered vertically in rectangle
        width = Inches(4.6) # Slightly less than rectangle width for margin
        height = Inches(2)  # Enough height for text
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.clear()  # Clear existing text
        text_frame.word_wrap = True  # Enable word wrap
        text_frame.vertical_anchor = 1  # Middle vertical alignment
        
        # Split text into lines
        lines = input.text.split('\n')
        logger.debug(f"Number of lines: {len(lines)}")
        logger.debug(f"Lines to add: {lines}")
        
        # Add each line as a separate paragraph
        for i, line in enumerate(lines):
            if line.strip():  # Only add non-empty lines
                p = text_frame.add_paragraph()
                p.text = line.strip()
                p.alignment = 1  # Center align the text
                
                # Format the text
                run = p.runs[0]
                if "Final Result:" in line:
                    run.font.size = Pt(32)  # Header size
                    run.font.bold = True
                else:
                    run.font.size = Pt(28)  # Value size
                    run.font.bold = True
                
                run.font.color.rgb = RGBColor(0, 0, 0)  # Black text
        prs.save(PPTX_FILENAME)
        await asyncio.sleep(5)
        
        # Reopen PowerPoint
        os.startfile(PPTX_FILENAME)
        await asyncio.sleep(10)
        # Reopen PowerPoint
        os.startfile(PPTX_FILENAME)
        await asyncio.sleep(10)
        
        logger.info(f"Text added successfully: {input.text}")
        return {
            "content": [
                TextContent(
                    type="text",
                    text=f"Text added successfully: {input.text}"
                )
            ]
        }
    except Exception as e:
        logger.error(f"Error in add_text_in_powerpoint: {str(e)}")
        logger.error(traceback.format_exc())
        return {
            "content": [
                TextContent(
                    type="text",
                    text=f"Error adding text: {str(e)}"
                )
            ]
        }

@mcp.tool()
async def send_gmail(input: SendGmailInput) -> dict:
    """Send an email with the specified content via Gmail"""
    try:
        logger.info(f"Calling send_gmail(content: {input.content[:50]}...)")
        
        # Retrieve Gmail credentials and recipient from .env
        gmail_address = os.getenv("GMAIL_ADDRESS")
        gmail_app_password = os.getenv("GMAIL_APP_PASSWORD")
        recipient_email = os.getenv("RECIPIENT_EMAIL")
        
        if not all([gmail_address, gmail_app_password, recipient_email]):
            error_msg = "Missing GMAIL_ADDRESS, GMAIL_APP_PASSWORD, or RECIPIENT_EMAIL in .env file"
            logger.error(error_msg)
            return {
                "content": [
                    TextContent(
                        type="text",
                        text=error_msg
                    )
                ]
            }
        
        # Validate Gmail address and recipient email format
        if not (gmail_address.endswith('@gmail.com') and '@' in recipient_email):
            error_msg = f"Invalid email format: GMAIL_ADDRESS={gmail_address}, RECIPIENT_EMAIL={recipient_email}"
            logger.error(error_msg)
            return {
                "content": [
                    TextContent(
                        type="text",
                        text=error_msg
                    )
                ]
            }
        
        # Create the email message
        msg = MIMEText(input.content)
        msg['Subject'] = 'Math Agent Result'
        msg['From'] = gmail_address
        msg['To'] = recipient_email
        
        # Connect to Gmail's SMTP server
        try:
            with smtplib.SMTP_SSL('smtp.gmail.com', 465) as server:
                logger.debug("Connecting to Gmail SMTP server (smtp.gmail.com:465)")
                server.login(gmail_address, gmail_app_password)
                logger.debug(f"Logged in as {gmail_address}")
                server.sendmail(gmail_address, recipient_email, msg.as_string())
                logger.info(f"Email sent successfully to {recipient_email}")
        except smtplib.SMTPAuthenticationError as e:
            error_msg = f"SMTP Authentication failed: {str(e)}. Ensure GMAIL_APP_PASSWORD is correct and 2-Step Verification is enabled."
            logger.error(error_msg)
            return {
                "content": [
                    TextContent(
                        type="text",
                        text=error_msg
                    )
                ]
            }
        except Exception as e:
            error_msg = f"Failed to send email: {str(e)}"
            logger.error(error_msg)
            return {
                "content": [
                    TextContent(
                        type="text",
                        text=error_msg
                    )
                ]
            }
        
        return {
            "content": [
                TextContent(
                    type="text",
                    text=f"Email sent successfully to {recipient_email}"
                )
            ]
        }
    except Exception as e:
        error_msg = f"Error in send_gmail: {str(e)}"
        logger.error(error_msg)
        logger.error(traceback.format_exc())
        return {
            "content": [
                TextContent(
                    type="text",
                    text=error_msg
                )
            ]
        }

# TOOLS WRAPPING FUNCTIONS FROM tools.py (with Pydantic validation)
@mcp.tool()
def t_number_list_to_sum(input: NumberListInput) -> NumberListOutput:
    """Sum numbers in a list"""
    logger.info(f"Calling t_number_list_to_sum with {len(input.numbers)} numbers")
    result = number_list_to_sum(input.numbers)
    return NumberListOutput(result=result)

@mcp.tool()
def t_calculate_difference(input: TwoNumberInput) -> TwoNumberOutput:
    """Difference between two numbers"""
    logger.info(f"Calling t_calculate_difference({input.a}, {input.b})")
    result = calculate_difference(input.a, input.b)
    return TwoNumberOutput(result=result)

@mcp.tool()
def t_number_list_to_product(input: NumberListInput) -> NumberListOutput:
    """Product of numbers in a list"""
    logger.info(f"Calling t_number_list_to_product with {len(input.numbers)} numbers")
    result = number_list_to_product(input.numbers)
    return NumberListOutput(result=result)

@mcp.tool()
def t_calculate_division(input: TwoNumberInput) -> TwoNumberOutput:
    """Division of two numbers"""
    logger.info(f"Calling t_calculate_division({input.a}, {input.b})")
    result = calculate_division(input.a, input.b)
    return TwoNumberOutput(result=result)

@mcp.tool()
def t_strings_to_chars_to_int(input: StringToCharsInput) -> StringToCharsOutput:
    """ASCII values of characters"""
    logger.info(f"Calling t_strings_to_chars_to_int('{input.text}')")
    ascii_values = local_strings_to_chars_to_int(input.text)
    return StringToCharsOutput(ascii_values=ascii_values)

@mcp.tool()
def t_int_list_to_exponential_values(input: ExponentialInput) -> ExponentialOutput:
    """Exponential of list elements"""
    logger.info(f"Calling t_int_list_to_exponential_values with {len(input.numbers)} numbers")
    values = int_list_to_exponential_values(input.numbers)
    return ExponentialOutput(values=values)

@mcp.tool()
def t_fibonacci_numbers(input: FibonacciInput) -> FibonacciOutput:
    """First n Fibonacci numbers"""
    logger.info(f"Calling t_fibonacci_numbers(n={input.n})")
    sequence = local_fibonacci_numbers(input.n)
    return FibonacciOutput(sequence=sequence)

@mcp.tool()
def t_calculate_factorial(input: FactorialInput) -> FactorialOutput:
    """List of factorials up to n-1"""
    logger.info(f"Calling t_calculate_factorial(n={input.n})")
    factorials = calculate_factorial(input.n)
    return FactorialOutput(factorials=factorials)

@mcp.tool()
def t_calculate_permutation(input: PermutationInput) -> PermutationOutput:
    """Permutation nPr"""
    logger.info(f"Calling t_calculate_permutation(n={input.n}, r={input.r})")
    result = calculate_permutation(input.n, input.r)
    return PermutationOutput(result=result)

@mcp.tool()
def t_calculate_combination(input: CombinationInput) -> CombinationOutput:
    """Combination nCr"""
    logger.info(f"Calling t_calculate_combination(n={input.n}, r={input.r})")
    result = calculate_combination(input.n, input.r)
    return CombinationOutput(result=result)

@mcp.tool()
def t_calculate_salary_for_id(input: EmployeeIdInput) -> SalaryOutput:
    """Salary by employee id"""
    logger.info(f"Calling t_calculate_salary_for_id(emp_id={input.emp_id})")
    salary = calculate_salary_for_id(input.emp_id)
    return SalaryOutput(salary=salary, found=salary is not None)

@mcp.tool()
def t_calculate_salary_for_name(input: EmployeeNameInput) -> SalaryOutput:
    """Salary by employee name"""
    logger.info(f"Calling t_calculate_salary_for_name(emp_name='{input.emp_name}')")
    salary = calculate_salary_for_name(input.emp_name)
    return SalaryOutput(salary=salary, found=salary is not None)

@mcp.tool()
def t_calculate_percentage(input: PercentageInput) -> PercentageOutput:
    """Calculate percentage of a number"""
    logger.info(f"Calling t_calculate_percentage(percent={input.percent}, number={input.number})")
    result = calculate_percentage(input.percent, input.number)
    return PercentageOutput(result=result)

@mcp.tool()
def fallback_reasoning(input: FallbackInput) -> FallbackOutput:
    """Fallback reasoning step when the agent is uncertain or a tool fails"""
    logger.info(f"Calling fallback_reasoning: {input.description}")
    message = f"Fallback invoked: {input.description}"
    return FallbackOutput(message=message)

if __name__ == "__main__":
    logger.info("Starting the MCP server")
    if len(sys.argv) > 1 and sys.argv[1] == "dev":
        logger.info("Running MCP server in dev mode without transport")
        mcp.run()  # Run without transport for dev server
    else:
        logger.info("Running MCP server with stdio transport")
        mcp.run(transport="stdio")  # Run with stdio for direct execution